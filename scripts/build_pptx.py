"""Build docs/overview.pptx — a PowerPoint port of docs/overview.html.

Usage:
    pip install python-pptx
    python scripts/build_pptx.py            # writes docs/overview.pptx
    python scripts/build_pptx.py --out X    # custom path
"""
from __future__ import annotations

import argparse
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

BG = RGBColor(0x10, 0x14, 0x1C)
FG = RGBColor(0xEE, 0xEE, 0xEE)
DIM = RGBColor(0xA8, 0xAE, 0xB8)
ACCENT = RGBColor(0x6A, 0xD1, 0xFF)
GOOD = RGBColor(0x50, 0xFA, 0x7B)
BAD = RGBColor(0xFF, 0x6B, 0x6B)
NEUTRAL = RGBColor(0xDD, 0xDD, 0xDD)

SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)  # 16:9, 13.33in × 7.5in
MARGIN = Emu(457200)                            # 0.5in
FONT = "Segoe UI"
MONO = "Consolas"


def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _textbox(slide, x: Emu, y: Emu, w: Emu, h: Emu):
    return slide.shapes.add_textbox(x, y, w, h).text_frame


def _set_run(run, *, font=FONT, size=18, bold=False, color=FG, mono=False) -> None:
    run.font.name = MONO if mono else font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _para(tf, text: str, *, size=18, bold=False, color=FG, mono=False, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color, mono=mono)
    return p


def _heading(slide, text: str, *, size=32, color=ACCENT) -> None:
    tf = _textbox(slide, MARGIN, MARGIN, SLIDE_W - 2 * MARGIN, Emu(700000))
    _para(tf, text, size=size, bold=True, color=color, first=True)


def _footer(slide, text: str = "bench-cleanser · v1.0.0") -> None:
    tf = _textbox(slide, MARGIN, SLIDE_H - Emu(380000), SLIDE_W - 2 * MARGIN, Emu(300000))
    _para(tf, text, size=10, color=DIM, align=PP_ALIGN.RIGHT, first=True)


def _table(slide, x: Emu, y: Emu, w: Emu, h: Emu, rows: list[list[tuple[str, RGBColor, bool]]], *, header=True, font_size=12):
    """Each cell is (text, color, bold)."""
    nrows, ncols = len(rows), len(rows[0])
    table = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    for r, row in enumerate(rows):
        for c, (text, color, bold) in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x18, 0x1F, 0x2C) if (header and r == 0) else BG
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(60000)
            tf.margin_right = Emu(60000)
            tf.margin_top = Emu(40000)
            tf.margin_bottom = Emu(40000)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            _set_run(run, size=font_size, bold=bold or (header and r == 0), color=color)
    return table


def _code_block(slide, x: Emu, y: Emu, w: Emu, h: Emu, lines: list[str], *, size=11) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0A, 0x10, 0x18)
    shape.line.color.rgb = RGBColor(0x25, 0x30, 0x40)
    tf = shape.text_frame
    tf.margin_left = Emu(180000)
    tf.margin_right = Emu(180000)
    tf.margin_top = Emu(120000)
    tf.margin_bottom = Emu(120000)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        _para(tf, line, size=size, color=ACCENT if i == 0 else FG, mono=True, first=(i == 0))


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    tf = _textbox(slide, MARGIN, Emu(2000000), SLIDE_W - 2 * MARGIN, Emu(900000))
    _para(tf, "bench-cleanser", size=60, bold=True, color=FG, align=PP_ALIGN.CENTER, first=True)
    tf = _textbox(slide, MARGIN, Emu(2900000), SLIDE_W - 2 * MARGIN, Emu(600000))
    _para(tf, "Deterministic contamination, fairness, and trajectory-leakage analysis", size=22, color=ACCENT, align=PP_ALIGN.CENTER, first=True)
    _para(tf, "for the SWE-bench family.", size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    tf = _textbox(slide, MARGIN, Emu(3700000), SLIDE_W - 2 * MARGIN, Emu(400000))
    _para(tf, "A research instrument for building, training on, and grading agentic coding LLMs.",
          size=14, color=DIM, align=PP_ALIGN.CENTER, first=True)
    tf = _textbox(slide, MARGIN, Emu(4500000), SLIDE_W - 2 * MARGIN, Emu(400000))
    _para(tf, "v1.0.0   ·   116 offline tests   ·   ruff + mypy clean   ·   Python 3.11 / 3.12   ·   MIT",
          size=13, color=GOOD, align=PP_ALIGN.CENTER, first=True)
    tf = _textbox(slide, MARGIN, Emu(5400000), SLIDE_W - 2 * MARGIN, Emu(300000))
    _para(tf, "Liao Zhu  ·  2026", size=11, color=DIM, align=PP_ALIGN.CENTER, first=True)


def slide_why(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "Why this exists")
    tf = _textbox(slide, MARGIN, Emu(1200000), SLIDE_W - 2 * MARGIN, Emu(900000))
    _para(tf, "SWE-bench leaderboards answer “did this row pass.”", size=18, first=True)
    _para(tf, "They don't answer “was passing this row a real measurement of capability.”", size=18, color=DIM)
    _para(tf, "Two failure modes that look identical on a leaderboard:", size=14, color=DIM)
    rows = [
        [("", DIM, True), ("Looks like", ACCENT, True), ("Actually is", ACCENT, True)],
        [("Task is broken", DIM, False), ("Model failed", FG, False), ("Tests demand undescribed behaviour", FG, False)],
        [("Agent cheated", DIM, False), ("Model passed", FG, False), ("Agent pip install-ed the fix", FG, False)],
        [("Honest pass", DIM, False), ("Model passed", FG, False), ("Model passed", FG, False)],
        [("Honest fail", DIM, False), ("Model failed", FG, False), ("Model failed", FG, False)],
    ]
    _table(slide, MARGIN, Emu(2700000), SLIDE_W - 2 * MARGIN, Emu(2200000), rows, font_size=14)
    tf = _textbox(slide, MARGIN, Emu(5100000), SLIDE_W - 2 * MARGIN, Emu(700000))
    _para(tf, "A single score conflates all four. bench-cleanser separates them, with citations, "
              "before anyone trains on the data or publishes against it.", size=14, color=DIM, first=True)
    _footer(slide)


def slide_two_axis(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "The two-axis model")
    col_w = (SLIDE_W - 3 * MARGIN) // 2
    tf = _textbox(slide, MARGIN, Emu(1300000), col_w, Emu(2800000))
    _para(tf, "Axis 1 — Task contamination", size=22, bold=True, color=ACCENT, first=True)
    _para(tf, "Is the benchmark item itself fair?", size=14, color=DIM)
    _para(tf, "•  Multi-label over 7 binary labels", size=14)
    _para(tf, "•  Severity = pure set membership (no thresholds, no floats)", size=14)
    _para(tf, "•  Inputs: problem text, gold patch, F2P/P2P tests, repo state", size=14)
    tf = _textbox(slide, MARGIN * 2 + col_w, Emu(1300000), col_w, Emu(2800000))
    _para(tf, "Axis 2 — Agent trajectory", size=22, bold=True, color=ACCENT, first=True)
    _para(tf, "How did this agent reach its result?", size=14, color=DIM)
    _para(tf, "•  One label per (task, agent)", size=14)
    _para(tf, "•  Heuristics + LLM, with cross-agent quorum upgrade", size=14)
    _para(tf, "•  Inputs: action trace, final patch, resolved flag", size=14)
    tf = _textbox(slide, MARGIN, Emu(4400000), SLIDE_W - 2 * MARGIN, Emu(1500000))
    _para(tf, "An APPROACH_LOCK task is broken whether or not any model attempts it.", size=14, first=True)
    _para(tf, "An agent that pip install-s the fix has cheated whether or not the task is clean.", size=14)
    _para(tf, "Stage 7 fuses both axes into a single fairness verdict.", size=14, bold=True, color=ACCENT)
    _footer(slide)


def slide_axis1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "Axis 1  ·  labels and severity")
    rows = [
        [("Label", ACCENT, True), ("Triggers when…", ACCENT, True), ("Severity", ACCENT, True)],
        [("APPROACH_LOCK", FG, True), ("F2P tests assert on a specific implementation, not behaviour.", FG, False), ("SEVERE", BAD, True)],
        [("OVER_TEST", FG, True), ("F2P tests assert on undescribed behaviour.", FG, False), ("SEVERE", BAD, True)],
        [("OVER_PATCH", FG, True), ("Gold patch modifies behaviour beyond the problem.", FG, False), ("MINOR / MODERATE", FG, False)],
        [("UNCLEAR_DESCRIPTION", FG, True), ("Multiple incompatible solutions are reasonable.", FG, False), ("MINOR", FG, False)],
        [("HIDDEN_CONTEXT", FG, True), ("Spec depends on info not in the problem text.", FG, False), ("MINOR", FG, False)],
        [("WEAK_COVERAGE", FG, True), ("F2P tests don't exercise patched code paths.", FG, False), ("MINOR", FG, False)],
        [("CLEAN", FG, True), ("None of the above.", FG, False), ("CLEAN", GOOD, True)],
    ]
    _table(slide, MARGIN, Emu(1200000), SLIDE_W - 2 * MARGIN, Emu(3500000), rows, font_size=12)
    _code_block(slide, MARGIN, Emu(4900000), SLIDE_W - 2 * MARGIN, Emu(1400000), [
        "SEVERE   ≜ APPROACH_LOCK ∈ L  OR  OVER_TEST ∈ L",
        "MODERATE ≜ OVER_PATCH ∈ L  AND  (HIDDEN_CONTEXT ∈ L  OR  UNCLEAR_DESCRIPTION ∈ L)",
        "MINOR    ≜ any other non-empty contamination subset",
        "CLEAN    ≜ L = ∅",
    ], size=12)
    _footer(slide)


def slide_axis2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "Axis 2  ·  per-agent trajectory labels")
    rows = [
        [("Outcome", ACCENT, True), ("Label", ACCENT, True), ("Pattern", ACCENT, True)],
        [("Passed", DIM, False), ("agent_passed_genuine", FG, True), ("Explore → hypothesise → patch → test.", FG, False)],
        [("", DIM, False), ("agent_passed_leak", FG, True), ("Final patch ≥ 0.90 similar to gold; direct file jumps.", FG, False)],
        [("", DIM, False), ("agent_passed_package_leak", FG, True), ("Agent pip install-ed the affected package.", FG, False)],
        [("", DIM, False), ("agent_passed_test_aware", FG, True), ("Referenced F2P test names before they were derivable.", FG, False)],
        [("", DIM, False), ("agent_passed_trained_hack", FG, True), ("Canonical fix on first try — memorised template.", FG, False)],
        [("Failed", DIM, False), ("agent_failed_completed_intent", FG, True), ("Addressed the brief; F2P tests rejected it → UNFAIR_FAILURE driver.", FG, False)],
        [("", DIM, False), ("agent_failed_no_intent", FG, True), ("Never engaged the problem. Skill gap.", FG, False)],
        [("Unknown", DIM, False), ("agent_unknown", FG, True), ("Trajectory truncated or malformed.", FG, False)],
    ]
    _table(slide, MARGIN, Emu(1200000), SLIDE_W - 2 * MARGIN, Emu(4800000), rows, font_size=12)
    _footer(slide)


def slide_fusion(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "Stage 7  ·  fairness verdict matrix")
    rows = [
        [("", ACCENT, True), ("passed_genuine", ACCENT, True), ("passed_leak·*", ACCENT, True),
         ("failed_completed", ACCENT, True), ("failed_no_intent", ACCENT, True), ("unknown", ACCENT, True)],
        [("CLEAN / MINOR", ACCENT, True), ("FAIR_PASS", GOOD, True), ("AGENT_CHEATED", BAD, True),
         ("FAIR_FAILURE", GOOD, True), ("AGENT_DISENGAGED", NEUTRAL, False), ("AMBIGUOUS_PASS\nINCONCLUSIVE", NEUTRAL, False)],
        [("MODERATE / SEVERE", ACCENT, True), ("CONTAMINATED_PASS", BAD, True), ("AGENT_CHEATED", BAD, True),
         ("UNFAIR_FAILURE", BAD, True), ("AGENT_DISENGAGED", NEUTRAL, False), ("INCONCLUSIVE", NEUTRAL, False)],
    ]
    _table(slide, MARGIN, Emu(1400000), SLIDE_W - 2 * MARGIN, Emu(2200000), rows, font_size=11)
    tf = _textbox(slide, MARGIN, Emu(4000000), SLIDE_W - 2 * MARGIN, Emu(1300000))
    _para(tf, "Red verdicts set invalidates_measurement = True — the single boolean a downstream consumer",
          size=14, first=True)
    _para(tf, "uses to drop a row from a leaderboard. Every verdict ships with reasoning and evidence: list[str].", size=14)
    _para(tf, " ", size=8)
    _para(tf, "Deterministic — no LLM call. Reproducible bit-for-bit from the persisted axis labels.",
          size=12, color=DIM)
    _footer(slide)


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "Architecture  ·  seven stages")
    lines = [
        "Stage 1     parse        diffs from gold patch + test patch              [deterministic]",
        "Stage 1.5   visit        shallow-clone repo · attach test source · AST   [deterministic]",
        "Stage 2     intent       LLM extracts core requirement + scope           [LLM]",
        "Stage 3     structural   astred-core / stdlib ast diff                   [deterministic]",
        "Stage 4A    patch ↔ intent                                                [LLM]",
        "Stage 4B    test  ↔ intent                                                [LLM]",
        "Stage 4C    cross-ref    overpatch–overtest coupling                     [deterministic]",
        "Stage 5     classify     dual-taxonomy LLM, heuristic-protected          [LLM]",
        "Stage 6     severity     set-membership bucket                           [deterministic]",
        "─────────────────────────────────────────────────────────────────────",
        "Stage 7     fusion       Axis 1 × Axis 2 → FairnessVerdict               [deterministic]",
    ]
    _code_block(slide, MARGIN, Emu(1200000), SLIDE_W - 2 * MARGIN, Emu(4200000), lines, size=11)
    tf = _textbox(slide, MARGIN, Emu(5600000), SLIDE_W - 2 * MARGIN, Emu(700000))
    _para(tf, "Every LLM stage emits a Pydantic BaseModel; severity, structural diff, cross-ref, "
              "and fusion are all deterministic. The package contains ZERO random.* calls.",
          size=13, color=DIM, first=True)
    _footer(slide)


def slide_determinism(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "Determinism contract")
    rows = [
        [("Guarantee", ACCENT, True), ("Concretely", ACCENT, True)],
        [("No floats in severity / fusion", FG, True),
         ("Severity = set membership. Fusion = (severity × labels × traj × resolved) → verdict.", FG, False)],
        [("No random", FG, True),
         ("LLM backoff jitter is (attempt mod 4) · 0.25 · base — deterministic.", FG, False)],
        [("Schema-enforced LLM I/O", FG, True),
         ("response_format={\"type\":\"json_schema\",\"strict\":true} with json_object fallback. "
          "Cache-before-validate so a corrupt payload stays inspectable.", FG, False)],
        [("Resumable runs", FG, True),
         ("--resume on by default; reports reloaded from disk are authoritative.", FG, False)],
        [("Frozen across model upgrades", FG, True),
         ("Severity rules depend only on label set. Swapping the LLM never changes the severity mapping.", FG, False)],
        [("Atomic writes", FG, True),
         ("Per-report JSON uses tempfile + os.replace; an interrupt never leaves half-files for --resume.", FG, False)],
        [("Cite-or-shut-up evidence", FG, True),
         ("A label without an evidence list is rejected by the classifier.", FG, False)],
    ]
    _table(slide, MARGIN, Emu(1200000), SLIDE_W - 2 * MARGIN, Emu(5000000), rows, font_size=12)
    _footer(slide)


def slide_trajectory(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "Trajectory infrastructure")
    tf = _textbox(slide, MARGIN, Emu(1100000), SLIDE_W - 2 * MARGIN, Emu(400000))
    _para(tf, "Axis 2 is where we deliberately invest beyond what a typical contamination tool ships.",
          size=14, color=DIM, first=True)
    rows = [
        [("Source", ACCENT, True), ("Loader", ACCENT, True), ("Notes", ACCENT, True)],
        [("Docent", FG, True), ("load_from_docent", FG, True), ("DQL → agent_runs; tool blocks mapped to ActionType.", FG, False)],
        [("HuggingFace", FG, True), ("load_from_huggingface", FG, True), ("Normalises SWE-bench agent dataset conventions.", FG, False)],
        [("JSONL", FG, True), ("load_from_jsonl", FG, True), ("One trajectory per line; tolerates malformed lines.", FG, False)],
        [("JSON dir", FG, True), ("load_from_json_dir", FG, True), ("One file per trajectory.", FG, False)],
    ]
    _table(slide, MARGIN, Emu(1700000), SLIDE_W - 2 * MARGIN, Emu(1700000), rows, font_size=12)
    tf = _textbox(slide, MARGIN, Emu(3600000), SLIDE_W - 2 * MARGIN, Emu(400000))
    _para(tf, "Layered classifier", size=18, bold=True, color=ACCENT, first=True)
    tf = _textbox(slide, MARGIN, Emu(4100000), SLIDE_W - 2 * MARGIN, Emu(2200000))
    _para(tf, "1.  Heuristics — normalised diff similarity (quote-aware comment strip), "
              "pip-install detection, F2P-name leakage.", size=13, first=True)
    _para(tf, "2.  LLM — strict-output Pydantic; heuristic signals embedded in prompt.", size=13)
    _para(tf, "3.  Cross-agent quorum — median pairwise similarity ≥ 0.85 AND ≥ 10 added lines "
              "→ upgrade to GOLD_PATCH_LEAK.", size=13)
    _para(tf, "      Outlier-tolerant; honest low-entropy convergence protected.", size=12, color=DIM)
    _footer(slide)


def slide_robustness(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "What we just hardened")
    col_w = (SLIDE_W - 3 * MARGIN) // 2
    cards = [
        ("Strict-mode schemas",
         "True json_schema strict:true with a _strictify_schema transformer "
         "(drops defaults, forces additionalProperties:false, inlines $defs). "
         "API rejects malformed candidates before they reach us."),
        ("Error sentinel removed",
         "Pipeline failures no longer masquerade as SEVERE. New pipeline_error "
         "field; severity stays CLEAN; Stage-7 fusion skips error rows. "
         "Aggregate stats no longer lie."),
        ("Concurrency throttled",
         "Trajectory analyzer respects max_concurrent_requests. Prevents "
         "rate-limit storms silently degrading LLM analysis to the heuristic fallback."),
        ("Cache stable across pydantic",
         "Structured cache key uses (model_name, schema_version), not the full "
         "schema text. Pydantic upgrades no longer wipe cached responses."),
        ("Heuristic union",
         "Three high-precision deterministic heuristics (pre-staged tests, "
         "0-REQUIRED mismatch, self-referential text) survive even if the LLM omits them."),
        ("Cross-agent quorum",
         "Median quorum + low-entropy gate. One-line bug fixes no longer get "
         "flagged as gold-patch leaks; one diverging agent no longer vetoes the upgrade."),
    ]
    card_h = Emu(1500000)
    gap = Emu(150000)
    for i, (title, body) in enumerate(cards):
        row, col = i // 2, i % 2
        x = MARGIN + col * (col_w + MARGIN)
        y = Emu(1200000) + row * (card_h + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x14, 0x1B, 0x28)
        shape.line.color.rgb = RGBColor(0x2A, 0x35, 0x48)
        tf = shape.text_frame
        tf.margin_left = Emu(180000)
        tf.margin_right = Emu(180000)
        tf.margin_top = Emu(140000)
        tf.margin_bottom = Emu(140000)
        tf.word_wrap = True
        _para(tf, title, size=15, bold=True, color=ACCENT, first=True)
        _para(tf, body, size=11, color=FG)
    _footer(slide)


def slide_status(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    _heading(slide, "Current status")
    rows = [
        [("Surface", ACCENT, True), ("State", ACCENT, True)],
        [("Source LOC", FG, True), ("~9.9k (33 modules)", FG, False)],
        [("Tests", FG, True), ("116 offline · no network · no Azure · no git clones", GOOD, False)],
        [("Lint", FG, True), ("ruff clean", GOOD, True)],
        [("Types", FG, True), ("mypy clean across all 33 modules", GOOD, True)],
        [("CI", FG, True), ("Python 3.11 + 3.12 matrix on Ubuntu", FG, False)],
        [("Benchmark coverage", FG, True), ("SWE-bench Verified · Pro · Live", FG, False)],
        [("Trajectory sources", FG, True), ("Docent · HuggingFace · JSONL · JSON dir", FG, False)],
        [("Console scripts", FG, True), ("bench-cleanser · bench-cleanser-trajectory · bench-cleanser-deep-dive", FG, False)],
    ]
    _table(slide, MARGIN, Emu(1200000), SLIDE_W - 2 * MARGIN, Emu(3300000), rows, font_size=12)
    tf = _textbox(slide, MARGIN, Emu(4700000), SLIDE_W - 2 * MARGIN, Emu(400000))
    _para(tf, "Honest caveats", size=16, bold=True, color=ACCENT, first=True)
    tf = _textbox(slide, MARGIN, Emu(5100000), SLIDE_W - 2 * MARGIN, Emu(1500000))
    _para(tf, "•  CLEAN means “no contamination signal on seven labels” — not “perfect benchmark item.”",
          size=12, color=DIM, first=True)
    _para(tf, "•  Axis-1 labels inherit LLM judgement noise; deterministic stages only amplify what upstream labels say.",
          size=12, color=DIM)
    _para(tf, "•  Cross-reference coupling is file-level, not function-level — known gap.",
          size=12, color=DIM)
    _para(tf, "•  Research instrument, not a metric. Makes human review tractable; does not replace it.",
          size=12, color=DIM)
    _footer(slide)


def slide_close(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG)
    tf = _textbox(slide, MARGIN, Emu(2000000), SLIDE_W - 2 * MARGIN, Emu(700000))
    _para(tf, "What it's for", size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, first=True)
    tf = _textbox(slide, MARGIN, Emu(2900000), SLIDE_W - 2 * MARGIN, Emu(1300000))
    _para(tf, "If you train on SWE-bench rows, every contaminated task teaches the wrong lesson.",
          size=20, align=PP_ALIGN.CENTER, first=True)
    _para(tf, "If you publish numbers against it, every contaminated row distorts the score.",
          size=20, align=PP_ALIGN.CENTER)
    tf = _textbox(slide, MARGIN, Emu(4400000), SLIDE_W - 2 * MARGIN, Emu(500000))
    _para(tf, "bench-cleanser is the tool you run before either.",
          size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, first=True)
    tf = _textbox(slide, MARGIN, Emu(5400000), SLIDE_W - 2 * MARGIN, Emu(400000))
    _para(tf, "github.com/v-liaozhu/bench-cleanser   ·   MIT   ·   v1.0.0",
          size=12, color=DIM, align=PP_ALIGN.CENTER, first=True)


def build(out_path: pathlib.Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for fn in (
        slide_title,
        slide_why,
        slide_two_axis,
        slide_axis1,
        slide_axis2,
        slide_fusion,
        slide_architecture,
        slide_determinism,
        slide_trajectory,
        slide_robustness,
        slide_status,
        slide_close,
    ):
        fn(prs)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote {out_path}  ({len(prs.slides)} slides)")


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--out", type=pathlib.Path, default=pathlib.Path("docs/overview.pptx"))
    args = ap.parse_args()
    build(args.out)


if __name__ == "__main__":
    main()
